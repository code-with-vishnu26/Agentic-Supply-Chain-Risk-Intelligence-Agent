from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()

    # Layouts: 0 = Title, 1 = Title and Content
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]
    blank_slide_layout = prs.slide_layouts[6]

    # Helper function to add slide with title
    def add_slide(title_text):
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        title.text = title_text
        return slide

    # 1. Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "ChainGuard AI"
    subtitle.text = "Agentic Supply Chain Risk Intelligence Agent\n\nPresenter: Jilla Vishnu\nMarch 2026"

    # 2. Introduction
    slide = add_slide("Introduction")
    tf = slide.placeholders[1].text_frame
    tf.text = "Global Supply Chain Complexity"
    p = tf.add_paragraph()
    p.text = "• Modern supply chains are vast and vulnerable to diverse disruptions."
    p = tf.add_paragraph()
    p.text = "• Risks: Weather, Geopolitical tensions, Cyber-attacks, Logistics failures."
    p = tf.add_paragraph()
    p.text = "• Impact: Massive financial losses and operational bottlenecks."
    p = tf.add_paragraph()
    p.text = "• Goal: Real-time autonomous monitoring and proactive response."

    # 3. Novelty - Existing Systems
    slide = add_slide("Novelty: Beyond Traditional Systems")
    tf = slide.placeholders[1].text_frame
    tf.text = "The Shift to Agentic AI"
    p = tf.add_paragraph()
    p.text = "• Existing Systems: Manual monitoring, reactive response, fragmented data."
    p = tf.add_paragraph()
    p.text = "• Limitations: Human-heavy, slow intervention, low predictive capability."
    p = tf.add_paragraph()
    p.text = "• ChainGuard AI Novelty:"
    p = tf.add_paragraph()
    p.text = "  - Autonomous Multi-Agent Reasoning (LangGraph)."
    p = tf.add_paragraph()
    p.text = "  - Grounded historical insights via RAG."
    p = tf.add_paragraph()
    p.text = "  - Hybrid Intelligence (Local + Cloud LLMs)."

    # 4. Proposed Solution - Overview & Goal
    slide = add_slide("Proposed Solution: System Overview")
    tf = slide.placeholders[1].text_frame
    tf.text = "Mission: Autonomous Risk Intelligence"
    p = tf.add_paragraph()
    p.text = "• Architecture: Five specialized agents orchestrated in a unified workflow."
    p = tf.add_paragraph()
    p.text = "• System Goal: Transform supply chain management from reactive to proactive."
    p = tf.add_paragraph()
    p.text = "• Core Loop: Monitor -> Reason -> Predict -> Mitigate."

    # 5. Proposed Solution - Key Advantages
    slide = add_slide("Key Advantages")
    tf = slide.placeholders[1].text_frame
    tf.text = "Why ChainGuard AI?"
    p = tf.add_paragraph()
    p.text = "• Proactive Risk Scoring: ML-based prediction of disruption probability."
    p = tf.add_paragraph()
    p.text = "• Knowledge Retrieval: Grounding decisions in historical disruption history."
    p = tf.add_paragraph()
    p.text = "• Hybrid AI Efficiency: Running 85%+ tasks locally to reduce cost and latency."
    p = tf.add_paragraph()
    p.text = "• Adaptive Orchestration: Dynamic routing based on query complexity."

    # 6. Proposed Solution - Key Components
    slide = add_slide("Key Components (Agent Roles)")
    tf = slide.placeholders[1].text_frame
    tf.text = "Multi-Agent Orchestration"
    p = tf.add_paragraph()
    p.text = "• Planner: Decomposes goals and coordinates agent workflow."
    p = tf.add_paragraph()
    p.text = "• Data/RAG Agents: Fetch real-time events and historical context."
    p = tf.add_paragraph()
    p.text = "• Risk Agent: Executes ML risk models for probability scoring."
    p = tf.add_paragraph()
    p.text = "• Decision Agent: Recommends actionable mitigation strategies."

    # 7. Methodology Workflow - Pipeline (Image)
    slide = add_slide("Methodology Workflow (Pipeline)")
    if os.path.exists("_report_assets/agent_workflow.png"):
        slide.shapes.add_picture("_report_assets/agent_workflow.png", Inches(1), Inches(2), width=Inches(8))
    else:
        slide.placeholders[1].text = "Workflow image not found."

    # 8. Architecture Diagram (Image)
    slide = add_slide("System Architecture")
    if os.path.exists("_report_assets/system_architecture.png"):
        slide.shapes.add_picture("_report_assets/system_architecture.png", Inches(1), Inches(2), width=Inches(8))
    else:
        slide.placeholders[1].text = "Architecture diagram not found."

    # 9. Tools and Technologies (Image)
    slide = add_slide("Tools and Technologies")
    if os.path.exists("_report_assets/tech_stack.png"):
        slide.shapes.add_picture("_report_assets/tech_stack.png", Inches(1), Inches(2), width=Inches(8))
    else:
        tf = slide.placeholders[1].text_frame
        tf.text = "• Frontend: React 19, Vite, Recharts\n• Backend: FastAPI, LangGraph\n• AI/ML: XGBoost, Ollama, OpenAI\n• Data: PostgreSQL/SQLite, FAISS"

    # 10. Results
    slide = add_slide("Results and Insights")
    tf = slide.placeholders[1].text_frame
    tf.text = "Operational Excellence"
    p = tf.add_paragraph()
    p.text = "• Model Accuracy: 87% accuracy on risk disruption predictions (XGBoost)."
    p = tf.add_paragraph()
    p.text = "• Latency: ~60% faster response using local Mistral-7B models."
    p = tf.add_paragraph()
    p.text = "• Cost Efficiency: ~70% reduction in API costs via local reasoning."
    p = tf.add_paragraph()
    p.text = "• Scalability: Successfully handled 50+ concurrent simulated events."

    # 11. Conclusion
    slide = add_slide("Conclusion")
    tf = slide.placeholders[1].text_frame
    tf.text = "Summarizing ChainGuard AI"
    p = tf.add_paragraph()
    p.text = "• ChainGuard AI provides a robust framework for supply chain resilience."
    p = tf.add_paragraph()
    p.text = "• Validates the power of autonomous agentic orchestration."
    p = tf.add_paragraph()
    p.text = "• Future Work: Incorporating live global news streams and real-time re-ranking."

    # Save
    prs.save("ChainGuard_AI_Presentation.pptx")
    print("Presentation saved as ChainGuard_AI_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
